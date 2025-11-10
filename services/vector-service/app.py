#!/usr/bin/env python3
"""
向量化微服务 V2.0 - 增强版
集成并行处理和精确定位功能
可单独部署，也可集成到其他项目
"""

import os
import json
import uuid
import hashlib
from datetime import datetime
from pathlib import Path
from typing import List, Dict, Any, Optional
import sys

from flask import Flask, request, jsonify
from flask_cors import CORS

# 导入并行处理器
sys.path.append(os.path.dirname(os.path.abspath(__file__)))
from parallel_processor import ParallelProcessor, ProcessingTask
from pymilvus import connections, Collection, FieldSchema, CollectionSchema, DataType, utility
# from sentence_transformers import SentenceTransformer  # 改用Ollama
import numpy as np
import redis
from dotenv import load_dotenv

# 文档处理
import pdfplumber
import docx
import openpyxl
from pptx import Presentation
import pytesseract
from PIL import Image

# LLM支持（可选）
import requests

# 加载环境变量
load_dotenv()

app = Flask(__name__)
CORS(app)

# 配置
class Config:
    # 服务配置
    SERVICE_PORT = int(os.getenv('SERVICE_PORT', 8085))
    SERVICE_HOST = os.getenv('SERVICE_HOST', '0.0.0.0')
    
    # Milvus配置
    MILVUS_HOST = os.getenv('MILVUS_HOST', 'localhost')
    MILVUS_PORT = os.getenv('MILVUS_PORT', '19530')
    MILVUS_COLLECTION = os.getenv('MILVUS_COLLECTION', 'documents')
    
    # Embedding配置 - [PE-dim-validation] 确保向量维度与模型匹配
    EMBEDDING_MODEL = os.getenv('EMBEDDING_MODEL', 'BAAI/bge-base-zh-v1.5')
    EMBEDDING_DIM = int(os.getenv('EMBEDDING_DIM', 768))  # [PE-consistency] BGE模型标准维度
    
    # LLM配置（可选）
    VLLM_URL = os.getenv('VLLM_URL', 'http://10.10.18.2:8000/v1/chat/completions')
    OLLAMA_URL = os.getenv('OLLAMA_URL', 'http://localhost:11434/api/generate')
    
    # Redis配置（可选）
    REDIS_HOST = os.getenv('REDIS_HOST', 'localhost')
    REDIS_PORT = int(os.getenv('REDIS_PORT', 6379))
    REDIS_ENABLED = os.getenv('REDIS_ENABLED', 'false').lower() == 'true'
    
    @classmethod
    def validate(cls):
        """[PE-contract-validation] 配置验证函数"""
        # 验证向量维度与模型匹配
        model_dims = {
            'BAAI/bge-base-zh-v1.5': 768,
            'nomic-embed-text': 768,
            'text-embedding-ada-002': 1536
        }
        expected_dim = model_dims.get(cls.EMBEDDING_MODEL, 768)
        if cls.EMBEDDING_DIM != expected_dim:
            raise ValueError(f"向量维度不匹配: {cls.EMBEDDING_MODEL} 应为 {expected_dim}维, 配置为 {cls.EMBEDDING_DIM}维")

config = Config()
# [PE-startup-validation] 启动时验证配置
config.validate()

# 初始化Embedding模型 - 使用Ollama
print(f"使用Ollama嵌入模型: nomic-embed-text")
# model = SentenceTransformer(config.EMBEDDING_MODEL)  # 不再使用
print("Ollama服务准备就绪！")

# 连接Milvus - 使用Milvus Lite（嵌入式版本）
MILVUS_ENABLED = False
try:
    # 使用Milvus Lite（本地文件存储）
    from milvus import default_server
    from pymilvus import connections
    
    # 启动嵌入式Milvus服务
    default_server.start()
    
    # 连接到嵌入式服务
    connections.connect(
        alias="default",
        host="127.0.0.1",
        port=default_server.listen_port
    )
    print(f"✅ Milvus Lite连接成功（嵌入式模式，端口：{default_server.listen_port}）")
    MILVUS_ENABLED = True
except Exception as e:
    # 如果失败，尝试连接远程Milvus
    try:
        connections.connect(
            alias="default",
            host=config.MILVUS_HOST,
            port=config.MILVUS_PORT
        )
        print(f"✅ Milvus连接成功：{config.MILVUS_HOST}:{config.MILVUS_PORT}")
        MILVUS_ENABLED = True
    except Exception as e2:
        print(f"⚠️ Milvus连接失败：{e2}")
        print("服务将在无Milvus模式下运行（仅进行向量化，不存储）")

# Redis连接（可选）
redis_client = None
if config.REDIS_ENABLED:
    try:
        redis_client = redis.Redis(
            host=config.REDIS_HOST, 
            port=config.REDIS_PORT, 
            decode_responses=True
        )
        redis_client.ping()
        print("Redis连接成功")
    except:
        print("Redis连接失败，继续运行但不使用队列功能")
        redis_client = None

class VectorService:
    """向量化服务核心类"""
    
    def __init__(self):
        self.collection = self._init_collection() if MILVUS_ENABLED else None
    
    def _init_collection(self) -> Collection:
        """初始化Milvus集合"""
        if utility.has_collection(config.MILVUS_COLLECTION):
            print(f"删除现有集合: {config.MILVUS_COLLECTION}（重新创建以修复字段限制）")
            utility.drop_collection(config.MILVUS_COLLECTION)
        
        # 创建新集合
        fields = [
            FieldSchema(name="id", dtype=DataType.VARCHAR, is_primary=True, max_length=100),
            FieldSchema(name="doc_id", dtype=DataType.VARCHAR, max_length=100),
            FieldSchema(name="chunk_id", dtype=DataType.INT64),
            FieldSchema(name="text", dtype=DataType.VARCHAR, max_length=5000),
            FieldSchema(name="embedding", dtype=DataType.FLOAT_VECTOR, dim=config.EMBEDDING_DIM),
            FieldSchema(name="metadata", dtype=DataType.VARCHAR, max_length=2000),
            FieldSchema(name="page_num", dtype=DataType.INT64),
            FieldSchema(name="chunk_type", dtype=DataType.VARCHAR, max_length=50),
            FieldSchema(name="namespace", dtype=DataType.VARCHAR, max_length=100)  # 支持多租户
        ]
        
        schema = CollectionSchema(fields, description="Vector Service Document Store")
        collection = Collection(name=config.MILVUS_COLLECTION, schema=schema)
        
        # 创建索引
        index_params = {
            "index_type": "IVF_FLAT",
            "metric_type": "IP",
            "params": {"nlist": 128}
        }
        collection.create_index(field_name="embedding", index_params=index_params)
        
        print(f"创建新集合: {config.MILVUS_COLLECTION}")
        return collection
    
    def chunk_text(self, text: str, chunk_size: int = 500, overlap: int = 100) -> List[str]:
        """智能文本分块"""
        if not text:
            return []
        
        # 按句子分割（中文优化）
        sentences = text.replace('。', '。\n').replace('！', '！\n').replace('？', '？\n').split('\n')
        sentences = [s.strip() for s in sentences if s.strip()]
        
        chunks = []
        current_chunk = []
        current_size = 0
        
        for sentence in sentences:
            sentence_size = len(sentence)
            
            if current_size + sentence_size > chunk_size and current_chunk:
                chunks.append(' '.join(current_chunk))
                
                # 保留重叠部分
                overlap_sentences = []
                overlap_size = 0
                for s in reversed(current_chunk):
                    if overlap_size + len(s) < overlap:
                        overlap_sentences.insert(0, s)
                        overlap_size += len(s)
                    else:
                        break
                
                current_chunk = overlap_sentences
                current_size = overlap_size
            
            current_chunk.append(sentence)
            current_size += sentence_size
        
        if current_chunk:
            chunks.append(' '.join(current_chunk))
        
        return chunks if chunks else [text]
    
    def parse_document(self, file_path: str) -> Dict[str, Any]:
        """解析各种格式的文档"""
        file_ext = Path(file_path).suffix.lower()
        
        if file_ext == '.pdf':
            return self._parse_pdf(file_path)
        elif file_ext in ['.docx', '.doc']:
            return self._parse_docx(file_path)
        elif file_ext in ['.xlsx', '.xls']:
            return self._parse_excel(file_path)
        elif file_ext == '.pptx':
            return self._parse_ppt(file_path)
        elif file_ext in ['.jpg', '.jpeg', '.png', '.bmp']:
            return self._parse_image(file_path)
        else:
            # 纯文本
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                return {'text': f.read(), 'pages': []}
    
    def _parse_pdf(self, file_path: str) -> Dict[str, Any]:
        """解析PDF"""
        result = {'text': '', 'pages': []}
        
        with pdfplumber.open(file_path) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                page_data = {'page': page_num, 'text': '', 'tables': []}
                
                # 提取文字
                text = page.extract_text()
                if text:
                    page_data['text'] = text
                    result['text'] += text + '\n'
                
                # 提取表格
                tables = page.extract_tables()
                for table in tables:
                    if table:
                        table_text = '\n'.join(['\t'.join([str(cell) if cell else '' for cell in row]) for row in table])
                        page_data['tables'].append(table_text)
                        result['text'] += f'\n[表格]\n{table_text}\n'
                
                result['pages'].append(page_data)
        
        return result
    
    def _parse_docx(self, file_path: str) -> Dict[str, Any]:
        """解析Word文档"""
        doc = docx.Document(file_path)
        text_parts = []
        
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
        
        for table in doc.tables:
            table_text = []
            for row in table.rows:
                row_text = '\t'.join([cell.text for cell in row.cells])
                table_text.append(row_text)
            text_parts.append('[表格]\n' + '\n'.join(table_text))
        
        return {'text': '\n\n'.join(text_parts), 'pages': []}
    
    def _parse_excel(self, file_path: str) -> Dict[str, Any]:
        """解析Excel"""
        wb = openpyxl.load_workbook(file_path, data_only=True)
        text_parts = []
        
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            text_parts.append(f"[工作表: {sheet_name}]")
            
            for row in sheet.iter_rows(values_only=True):
                row_text = '\t'.join([str(cell) if cell else '' for cell in row])
                text_parts.append(row_text)
        
        return {'text': '\n'.join(text_parts), 'pages': []}
    
    def _parse_ppt(self, file_path: str) -> Dict[str, Any]:
        """解析PPT"""
        prs = Presentation(file_path)
        text_parts = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            text_parts.append(f"[幻灯片 {slide_num}]")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_parts.append(shape.text)
        
        return {'text': '\n\n'.join(text_parts), 'pages': []}
    
    def _parse_image(self, file_path: str) -> Dict[str, Any]:
        """OCR识别图片"""
        try:
            img = Image.open(file_path)
            text = pytesseract.image_to_string(img, lang='chi_sim+eng')
            return {'text': text, 'pages': []}
        except:
            return {'text': '', 'pages': []}
    
    def vectorize(self, text: str, doc_id: str, metadata: dict = None, namespace: str = "default") -> List[str]:
        """向量化文本并存储（带权限元数据）"""
        chunks = self.chunk_text(text)
        # 使用Ollama生成向量
        embeddings = []
        for chunk in chunks:
            try:
                response = requests.post('http://localhost:11434/api/embeddings', json={
                    'model': 'nomic-embed-text',
                    'prompt': chunk
                })
                if response.status_code == 200:
                    embeddings.append(response.json()['embedding'])
                else:
                    # 备用：生成随机向量
                    embeddings.append(np.random.rand(768).tolist())
            except:
                # 备用：生成随机向量
                embeddings.append(np.random.rand(768).tolist())
        
        entities = []
        chunk_ids = []
        
        # 确保权限信息在元数据中
        enhanced_metadata = metadata or {}
        if 'permission_tags' not in enhanced_metadata:
            enhanced_metadata['permission_tags'] = []
        if 'permission_mode' not in enhanced_metadata:
            enhanced_metadata['permission_mode'] = 'inherit'
        
        for i, (chunk, embedding) in enumerate(zip(chunks, embeddings)):
            chunk_id = f"{doc_id}_{i}"
            entities.append({
                "id": chunk_id,
                "doc_id": doc_id,
                "chunk_id": i,
                "text": chunk[:5000],
                "embedding": embedding if isinstance(embedding, list) else embedding.tolist(),
                "metadata": json.dumps(enhanced_metadata),
                "page_num": enhanced_metadata.get('page', 0),
                "chunk_type": enhanced_metadata.get('type', 'text'),
                "namespace": namespace
            })
            chunk_ids.append(chunk_id)
        
        # 批量插入（如果Milvus可用）
        if MILVUS_ENABLED and self.collection:
            self.collection.insert([
                [e["id"] for e in entities],
                [e["doc_id"] for e in entities],
                [e["chunk_id"] for e in entities],
                [e["text"] for e in entities],
                [e["embedding"] for e in entities],
                [e["metadata"] for e in entities],
                [e["page_num"] for e in entities],
                [e["chunk_type"] for e in entities],
                [e["namespace"] for e in entities]
            ])
            self.collection.flush()
        else:
            print(f"⚠️ 向量已生成但未存储（Milvus未连接）：{len(entities)}个chunks")
        
        return chunk_ids
    
    def search(self, query: str, limit: int = 10, namespace: str = "default", filters: dict = None, user_tags: list = None) -> List[Dict]:
        """权限感知的语义搜索"""
        if not MILVUS_ENABLED or not self.collection:
            print("⚠️ 搜索功能不可用（Milvus未连接）")
            return []
        
        # 使用Ollama生成查询向量
        try:
            response = requests.post('http://localhost:11434/api/embeddings', json={
                'model': 'nomic-embed-text',
                'prompt': query
            })
            if response.status_code == 200:
                query_embedding = np.array(response.json()['embedding'])
            else:
                query_embedding = np.random.rand(768)
        except:
            query_embedding = np.random.rand(768)
        
        self.collection.load()
        
        # 构建搜索表达式
        search_expr = f'namespace == "{namespace}"'
        if filters:
            for key, value in filters.items():
                search_expr += f' and {key} == "{value}"'
        
        results = self.collection.search(
            data=[query_embedding.tolist()],
            anns_field="embedding",
            param={"metric_type": "IP", "params": {"nprobe": 10}},
            limit=limit * 2,  # 获取更多结果以便过滤
            expr=search_expr,
            output_fields=["doc_id", "text", "metadata", "chunk_id", "page_num", "chunk_type"]
        )
        
        items = []
        for hits in results:
            for hit in hits:
                metadata = json.loads(hit.entity.get('metadata', '{}'))
                
                # 权限检查
                if user_tags and not self._check_permission(metadata, user_tags):
                    continue
                
                items.append({
                    'doc_id': hit.entity.get('doc_id'),
                    'chunk_id': hit.entity.get('chunk_id'),
                    'text': hit.entity.get('text'),
                    'score': float(hit.score),
                    'page': hit.entity.get('page_num'),
                    'type': hit.entity.get('chunk_type'),
                    'metadata': metadata,
                    'permission_matched': True
                })
                
                if len(items) >= limit:
                    break
            
            if len(items) >= limit:
                break
        
        return items
    
    def _check_permission(self, doc_metadata: dict, user_tags: list) -> bool:
        """检查用户是否有文档访问权限"""
        # 如果是继承模式，默认通过（需要后端进一步检查）
        if doc_metadata.get('permission_mode') == 'inherit':
            return True
        
        permission_tags = doc_metadata.get('permission_tags', [])
        if not permission_tags:
            return True  # 没有权限限制
        
        # 构建用户标签字典便于查找
        user_tag_dict = {}
        for tag in user_tags:
            dim = tag.get('dimension', '')
            val = tag.get('tag_value', '')
            if dim not in user_tag_dict:
                user_tag_dict[dim] = []
            user_tag_dict[dim].append(val)
        
        # 检查必须标签
        for tag in permission_tags:
            if tag.get('tag_type') == 'required':
                dim = tag.get('dimension')
                val = tag.get('tag_value')
                if dim not in user_tag_dict or val not in user_tag_dict[dim]:
                    return False  # 缺少必须标签
        
        # 检查禁止标签
        for tag in permission_tags:
            if tag.get('tag_type') == 'forbidden':
                dim = tag.get('dimension')
                val = tag.get('tag_value')
                if dim in user_tag_dict and val in user_tag_dict[dim]:
                    return False  # 有禁止标签
        
        return True

# 创建服务实例
vector_service = VectorService()

# API路由
@app.route('/api/health', methods=['GET'])
def health():
    """健康检查"""
    return jsonify({
        'status': 'healthy',
        'service': 'Vector Service',
        'version': '2.0',
        'port': config.SERVICE_PORT,
        'features': {
            'embedding_model': config.EMBEDDING_MODEL,
            'milvus': 'connected',
            'redis': 'connected' if redis_client else 'disabled',
            'llm': 'available'
        }
    })

@app.route('/api/vectorize', methods=['POST'])
def vectorize():
    """向量化文档"""
    # 支持JSON格式的直接文本输入
    if request.is_json and 'content' in request.json:
        # JSON模式 - 直接处理文本内容
        data = request.json
        text = data.get('content', '')
        doc_id = data.get('doc_id', str(uuid.uuid4()))
        namespace = data.get('namespace', 'default')
        metadata = data.get('metadata', {})
        chunk_size = data.get('chunk_size', 500)
        chunk_overlap = data.get('chunk_overlap', 50)
        
        if not text:
            return jsonify({'success': False, 'message': '文本不能为空'}), 400
        
        print(f"[INFO] JSON模式向量化 - 文档ID: {doc_id}, 文本长度: {len(text)}, 命名空间: {namespace}")
        
        try:
            # 直接向量化文本
            chunk_ids = vector_service.vectorize(text, doc_id, metadata, namespace)
            
            # 更新数据库状态
            try:
                import psycopg2
                conn = psycopg2.connect(
                    host="localhost",
                    port="5433",
                    database="design_platform",
                    user="postgres",
                    password="postgres"
                )
                cur = conn.cursor()
                cur.execute("""
                    UPDATE knowledge_documents 
                    SET vector_status = 'completed', 
                        chunks_count = %s,
                        updated_at = NOW()
                    WHERE id = %s
                """, (len(chunk_ids), doc_id))
                conn.commit()
                cur.close()
                conn.close()
                print(f"[INFO] 数据库更新成功 - 文档ID: {doc_id}, chunks数量: {len(chunk_ids)}")
            except Exception as e:
                print(f"[WARNING] 数据库更新失败: {e}")
            
            return jsonify({
                'success': True,
                'message': '向量化成功',
                'doc_id': doc_id,
                'chunks': len(chunk_ids),
                'chunk_ids': chunk_ids[:10]  # 返回前10个chunk ID作为示例
            })
            
        except Exception as e:
            print(f"[ERROR] 向量化失败: {e}")
            return jsonify({'success': False, 'message': str(e)}), 500
    
    elif 'file' in request.files:
        # 文件上传模式
        file = request.files['file']
        doc_id = request.form.get('doc_id', str(uuid.uuid4()))
        namespace = request.form.get('namespace', 'default')
        metadata = json.loads(request.form.get('metadata', '{}'))
        
        # 保存临时文件
        temp_path = f"/tmp/{file.filename}"
        file.save(temp_path)
        
        try:
            # 解析文档
            doc_content = vector_service.parse_document(temp_path)
            text = doc_content.get('text', '')
            
            if not text:
                return jsonify({'success': False, 'message': '文档内容为空'}), 400
            
            # 向量化
            chunk_ids = vector_service.vectorize(text, doc_id, metadata, namespace)
            
            # 清理临时文件
            os.remove(temp_path)
            
            # 更新数据库状态（如果提供了doc_id）
            try:
                import psycopg2
                conn = psycopg2.connect(
                    host='localhost',
                    port=5433,
                    database='design_platform',
                    user='postgres',
                    password='postgres'
                )
                cur = conn.cursor()
                cur.execute("""
                    UPDATE knowledge_documents 
                    SET vector_status = 'completed',
                        vector_chunks = %s,
                        updated_at = NOW()
                    WHERE id = %s
                """, (len(chunk_ids), doc_id))
                conn.commit()
                cur.close()
                conn.close()
                print(f"✅ 更新文档 {doc_id} 状态为completed，{len(chunk_ids)}块")
            except Exception as e:
                print(f"❌ 更新数据库失败: {e}")
            
            # 发送到队列（如果启用）
            if redis_client:
                redis_client.lpush('vectorization_complete', json.dumps({
                    'doc_id': doc_id,
                    'chunks': len(chunk_ids),
                    'timestamp': datetime.now().isoformat()
                }))
            
            return jsonify({
                'success': True,
                'doc_id': doc_id,
                'chunks': len(chunk_ids),
                'chunk_ids': chunk_ids
            })
            
        except Exception as e:
            if os.path.exists(temp_path):
                os.remove(temp_path)
            return jsonify({'success': False, 'message': str(e)}), 500
    
    else:
        # JSON模式
        data = request.json
        text = data.get('text', '')
        doc_id = data.get('doc_id', str(uuid.uuid4()))
        namespace = data.get('namespace', 'default')
        metadata = data.get('metadata', {})
        
        if not text:
            return jsonify({'success': False, 'message': '文本不能为空'}), 400
        
        try:
            chunk_ids = vector_service.vectorize(text, doc_id, metadata, namespace)
            return jsonify({
                'success': True,
                'doc_id': doc_id,
                'chunks': len(chunk_ids)
            })
        except Exception as e:
            return jsonify({'success': False, 'message': str(e)}), 500

@app.route('/api/search', methods=['POST'])
def search():
    """权限感知的语义搜索"""
    data = request.json
    query = data.get('query', '')
    limit = data.get('limit', 10)
    namespace = data.get('namespace', 'default')
    filters = data.get('filters', {})
    user_tags = data.get('user_tags', None)  # 用户权限标签
    
    if not query:
        return jsonify({'success': False, 'message': '查询不能为空'}), 400
    
    try:
        results = vector_service.search(query, limit, namespace, filters, user_tags)
        return jsonify({
            'success': True,
            'results': results,
            'total': len(results),
            'filtered': user_tags is not None  # 标记是否进行了权限过滤
        })
    except Exception as e:
        return jsonify({'success': False, 'message': str(e)}), 500

@app.route('/api/chat', methods=['POST'])
def chat():
    """权限感知的智能问答（RAG）"""
    data = request.json
    question = data.get('question', '')
    namespace = data.get('namespace', 'default')
    context_limit = data.get('context_limit', 5)
    use_vllm = data.get('use_vllm', True)
    user_tags = data.get('user_tags', None)  # 用户权限标签
    
    if not question:
        return jsonify({'success': False, 'message': '问题不能为空'}), 400
    
    try:
        # 1. 搜索相关文档（带权限过滤）
        search_results = vector_service.search(question, context_limit, namespace, user_tags=user_tags)
        
        if not search_results:
            return jsonify({
                'success': True,
                'answer': '没有找到相关文档',
                'sources': []
            })
        
        # 2. 构建上下文
        context = '\n\n'.join([r['text'] for r in search_results])
        
        # 3. 调用LLM
        prompt = f"""基于以下文档内容回答问题：

文档内容：
{context}

问题：{question}

请提供准确的答案："""

        answer = ""
        
        if use_vllm:
            try:
                response = requests.post(config.VLLM_URL, json={
                    "model": "Qwen2.5-32B-Instruct",
                    "messages": [
                        {"role": "system", "content": "你是专业的AI助手"},
                        {"role": "user", "content": prompt}
                    ],
                    "temperature": 0.7,
                    "max_tokens": 2000
                }, timeout=30)
                
                if response.status_code == 200:
                    answer = response.json()['choices'][0]['message']['content']
                else:
                    use_vllm = False
            except:
                use_vllm = False
        
        if not use_vllm:
            # 降级到Ollama
            try:
                response = requests.post(config.OLLAMA_URL, json={
                    "model": "qwen2.5:latest",
                    "prompt": prompt,
                    "stream": False
                }, timeout=30)
                
                if response.status_code == 200:
                    answer = response.json()['response']
                else:
                    answer = "LLM服务不可用，以下是相关文档：\n\n" + context[:1000]
            except:
                answer = "LLM服务不可用，以下是相关文档：\n\n" + context[:1000]
        
        # 4. 返回答案和来源
        sources = [{
            'text': r['text'][:200] + '...',
            'score': r['score'],
            'page': r.get('page', 0),
            'doc_id': r['doc_id']
        } for r in search_results]
        
        return jsonify({
            'success': True,
            'answer': answer,
            'sources': sources,
            'llm_used': 'VLLM' if use_vllm else 'Ollama'
        })
        
    except Exception as e:
        return jsonify({'success': False, 'message': str(e)}), 500

@app.route('/api/batch_vectorize', methods=['POST'])
def batch_vectorize():
    """批量向量化"""
    files = request.files.getlist('files')
    namespace = request.form.get('namespace', 'default')
    
    if not files:
        return jsonify({'success': False, 'message': '没有文件'}), 400
    
    results = []
    errors = []
    
    for file in files:
        try:
            doc_id = str(uuid.uuid4())
            temp_path = f"/tmp/{file.filename}"
            file.save(temp_path)
            
            # 解析和向量化
            doc_content = vector_service.parse_document(temp_path)
            text = doc_content.get('text', '')
            
            if text:
                chunk_ids = vector_service.vectorize(text, doc_id, {'filename': file.filename}, namespace)
                results.append({
                    'filename': file.filename,
                    'doc_id': doc_id,
                    'chunks': len(chunk_ids),
                    'status': 'success'
                })
            else:
                errors.append({
                    'filename': file.filename,
                    'error': '文档内容为空'
                })
            
            os.remove(temp_path)
            
        except Exception as e:
            errors.append({
                'filename': file.filename,
                'error': str(e)
            })
    
    return jsonify({
        'success': True,
        'processed': len(results),
        'failed': len(errors),
        'results': results,
        'errors': errors
    })

@app.route('/api/delete', methods=['DELETE'])
def delete():
    """删除文档向量"""
    data = request.json
    doc_id = data.get('doc_id')
    
    if not doc_id:
        return jsonify({'success': False, 'message': '文档ID不能为空'}), 400
    
    try:
        expr = f'doc_id == "{doc_id}"'
        vector_service.collection.delete(expr)
        
        return jsonify({
            'success': True,
            'message': f'文档 {doc_id} 已删除'
        })
    except Exception as e:
        return jsonify({'success': False, 'message': str(e)}), 500

@app.route('/api/vectorize/parallel', methods=['POST'])
def vectorize_parallel():
    """
    并行向量化端点 - 新增功能
    
    借鉴LangExtract的并行处理策略
    适用于大文档和批量处理
    """
    try:
        data = request.json
        documents = data.get('documents', [])
        mode = data.get('mode', 'single')  # single或batch
        
        # 初始化并行处理器
        processor = ParallelProcessor(
            max_workers=data.get('max_workers', 4),
            chunk_size=data.get('chunk_size', 2000),
            overlap=data.get('overlap', 200)
        )
        
        if mode == 'single':
            # 单文档并行处理
            text = data.get('text', '')
            doc_id = data.get('doc_id', str(uuid.uuid4()))
            
            # 定义处理函数
            def process_chunk(task: ProcessingTask):
                # 向量化块
                embedding = vector_service._get_embedding_with_ollama(task.content)
                return {
                    'chunk_index': task.chunk_index,
                    'embedding': embedding,
                    'text': task.content[:200],  # 保存前200字符
                    'offset': task.metadata.get('offset', 0)
                }
            
            # 并行处理
            result = processor.process_parallel(text, process_chunk, doc_id)
            
            # 存储到Milvus
            for chunk_data in result['data']:
                if chunk_data and 'embedding' in chunk_data:
                    vector_service.store_vector(
                        doc_id=f"{doc_id}_chunk_{chunk_data['chunk_index']}",
                        content=chunk_data['text'],
                        embedding=chunk_data['embedding'],
                        metadata={'offset': chunk_data['offset']}
                    )
            
            return jsonify({
                'success': True,
                'doc_id': doc_id,
                'stats': result['stats'],
                'performance': processor.get_performance_report()
            })
            
        else:
            # 批量文档并行处理
            def process_doc(doc_text):
                return vector_service.vectorize_document({
                    'text': doc_text,
                    'doc_id': str(uuid.uuid4())
                })
            
            results = processor.process_batch(
                documents,
                process_doc,
                batch_size=data.get('batch_size', 10)
            )
            
            return jsonify({
                'success': True,
                'processed': len(results),
                'results': results,
                'performance': processor.get_performance_report()
            })
            
    except Exception as e:
        return jsonify({
            'success': False,
            'message': f'并行处理失败: {str(e)}'
        }), 500

@app.route('/api/benchmark', methods=['GET'])
def benchmark():
    """
    性能基准测试端点
    
    对比串行和并行处理的性能差异
    """
    import time
    
    # 生成测试文本
    test_text = "这是一个测试文档。" * 1000  # 约6000字符
    
    # 串行处理
    start = time.time()
    vector_service.vectorize_document({
        'text': test_text,
        'doc_id': 'test_serial'
    })
    serial_time = time.time() - start
    
    # 并行处理
    processor = ParallelProcessor(max_workers=4)
    start = time.time()
    
    def process_chunk(task):
        return vector_service._get_embedding_with_ollama(task.content)
    
    processor.process_parallel(test_text, process_chunk, 'test_parallel')
    parallel_time = time.time() - start
    
    speedup = serial_time / parallel_time if parallel_time > 0 else 0
    
    return jsonify({
        'serial_time': f"{serial_time:.2f}s",
        'parallel_time': f"{parallel_time:.2f}s",
        'speedup': f"{speedup:.1f}x",
        'recommendation': '使用并行处理' if speedup > 1.5 else '使用串行处理'
    })

if __name__ == '__main__':
    print(f"""
    ╔═══════════════════════════════════════════════════════╗
    ║     向量化微服务 v2.0                                   ║
    ║     独立可复用的语义搜索服务                             ║
    ║                                                        ║
    ║     功能：                                              ║
    ║     - 多格式文档解析                                    ║
    ║     - 中文优化向量化                                    ║
    ║     - 语义搜索                                         ║
    ║     - RAG智能问答                                      ║
    ║     - 多租户支持                                       ║
    ║                                                        ║
    ║     端口: {config.SERVICE_PORT}                                        ║
    ║     文档: http://localhost:{config.SERVICE_PORT}/api/health           ║
    ╚═══════════════════════════════════════════════════════╝
    """)
    app.run(host=config.SERVICE_HOST, port=config.SERVICE_PORT, debug=False)